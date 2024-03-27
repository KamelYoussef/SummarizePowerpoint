import streamlit as st
from pptx import Presentation
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import OllamaEmbeddings
from langchain.callbacks.manager import CallbackManager
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
from langchain_community.llms import Ollama
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template


def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)


def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks


def get_vectorstore(text_chunks):
    embeddings = OllamaEmbeddings(base_url = "http://localhost:11434", model="stablelm-zephyr")
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore


def get_conversation_chain(vectorstore):
    llm = Ollama(
        model="stablelm-zephyr",
        base_url="http://localhost:11434",
        verbose=True,
        callback_manager=CallbackManager([StreamingStdOutCallbackHandler()])
    )
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain


def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)



def main():
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    user_question = st.chat_input("Ask a question about your documents ...")
    if user_question:
        handle_userinput(user_question)

    with st.sidebar:
        st.subheader("Ask your documents")
        uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"], accept_multiple_files=False)
        if st.button("Process the documents"):
            with st.spinner("In Progress"):
                text = extract_text_from_pptx(uploaded_file)
                chunks = get_text_chunks(text)
                vectorstore = get_vectorstore(chunks)
                st.session_state.conversation = get_conversation_chain(vectorstore)


if __name__ == "__main__":
    main()
